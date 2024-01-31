from flask import Flask, jsonify, request
from openai import OpenAI
from pptx import Presentation
import os
from flask_cors import CORS

# Initialize OpenAI client with API key
client = OpenAI(api_key=os.environ.get("OPENAI_API_KEY"))

app = Flask(__name__)
CORS(app, resources={r"/generate_content": {"origins": "*"}})

# Define the base directory for shared data
SHARED_DATA_DIR = '/path/in/container'

@app.route('/generate_content', methods=['POST'])
def generate_content():
    if request.is_json:
        try:
            # ... [GPT-3 content generation code] ...

            # Extract the generated content from the response
            generated_content = completion.choices[0].message.content

            # Create a PowerPoint presentation
            ppt = Presentation()
            slide = ppt.slides.add_slide(ppt.slide_layouts[1])

            # Add content to the slide
            title = slide.shapes.title
            content = slide.placeholders[1]

            title.text = "Generated Content"
            content.text = generated_content

            # Define the path to save the presentation in the shared volume
            ppt_file = os.path.join(SHARED_DATA_DIR, 'generated_presentation.pptx')

            # Save the presentation
            ppt.save(ppt_file)

            # Return the path to the saved presentation
            return jsonify({'generated_content': 'Presentation created successfully', 'file': ppt_file})

        except Exception as e:
            return jsonify({'error': str(e)}), 500
    else:
        return jsonify({'error': 'Invalid Content-Type'}), 400

# ... [Flask app run configuration] ...
if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)
